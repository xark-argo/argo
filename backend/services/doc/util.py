import hashlib
import logging
import os.path
import random
import re
import zipfile
from collections.abc import Iterator
from pathlib import Path

import cchardet
import fitz

# import sentence_transformers
import xlrd
from bs4 import BeautifulSoup
from langchain_community.document_loaders import (
    BSHTMLLoader,
    CSVLoader,
    Docx2txtLoader,
    OutlookMessageLoader,
    TextLoader,
    UnstructuredEPubLoader,
    UnstructuredRSTLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredXMLLoader,
)
from langchain_core.document_loaders import BaseLoader
from langchain_core.documents import Document
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openpyxl import load_workbook
from pptx import Presentation

from core.callback_handler.logging_out_callback_handler import (
    LoggingOutCallbackHandler,
)
from core.model_providers import model_provider_manager
from core.model_providers.ollama.ollama_api import (
    ollama_model_is_generation,
)
from models.bot import BotModelConfig


class DocxHtmlLoader(BaseLoader):
    def __init__(self, file_path):
        self.file_path = file_path
        self.reg = re.compile(r"Generated Document(.*)------=mhtDocumentPart--", re.DOTALL)

    def extract_content(self):
        with zipfile.ZipFile(self.file_path, "r") as docx_zip:
            file_list = docx_zip.namelist()
            mht_files = [f for f in file_list if f.endswith(".mht") or f.endswith(".html")]
            if not mht_files:
                logging.info("mht files not found.")
                return None
            with docx_zip.open(mht_files[0]) as mht_file:
                content = mht_file.read().decode("utf-8")
                soup = BeautifulSoup(content, "html.parser")
                full_text = soup.get_text(separator="\n", strip=True)
                return self.reg.findall(full_text)[0]

    def lazy_load(self) -> Iterator[Document]:
        metadata = {"source": str(self.file_path)}
        try:
            context = self.extract_content()
            yield Document(page_content=context, metadata=metadata)
        except Exception as ex:
            logging.exception(f"Error extracting content from file: {self.file_path}.")
        yield Document(page_content="", metadata=metadata)


class PdfLoader(BaseLoader):
    def __init__(self, file_path):
        self.file_path = file_path

    def extract_content(self):
        try:
            if not os.path.exists(self.file_path):
                raise FileNotFoundError(f"pdf file {self.file_path} not exists")
            with fitz.open(self.file_path) as pdf_document:
                full_text = ""
                for page in pdf_document:
                    full_text += page.get_text()
                return full_text.strip()
        except Exception as ex:
            logging.exception("pdf extract fail")
            return None

    def lazy_load(self) -> Iterator[Document]:
        metadata = {"source": str(self.file_path)}
        try:
            context = self.extract_content()
            if context:
                yield Document(page_content=context, metadata=metadata)
            else:
                logging.info("no content extracted.")
                yield Document(page_content="", metadata=metadata)
        except Exception as ex:
            logging.exception(f"Error extracting content from file: {self.file_path}.")
            yield Document(page_content="", metadata=metadata)


class MarkdownLoader(BaseLoader):
    def __init__(self, file_path: str):
        self.file_path = file_path

    def extract_content(self):
        try:
            if not os.path.exists(self.file_path):
                raise FileNotFoundError(f"Markdown file {self.file_path} does not exist.")
            full_text = Path(self.file_path).read_text(encoding="utf-8")
            return full_text.strip()
        except Exception as ex:
            logging.exception("Markdown extraction failed")
            return None

    def lazy_load(self) -> Iterator[Document]:
        metadata = {"source": str(self.file_path)}
        try:
            context = self.extract_content()
            if context:
                yield Document(page_content=context, metadata=metadata)
            else:
                logging.info("No content extracted from Markdown file.")
                yield Document(page_content="", metadata=metadata)
        except Exception as ex:
            logging.exception(f"Error extracting content from file: {self.file_path}.")
            yield Document(page_content="", metadata=metadata)


class PPTLoader(BaseLoader):
    def __init__(self, file_path: str):
        self.file_path = file_path

    def extract_content(self):
        try:
            if not os.path.exists(self.file_path):
                raise FileNotFoundError(f"pptx file {self.file_path} does not exist.")

            presentation = Presentation(self.file_path)
            full_text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        full_text += shape.text.strip() + "\n"
            return full_text.strip()
        except Exception as ex:
            logging.exception("PPTX extraction failed.")
            return None

    def lazy_load(self) -> Iterator[Document]:
        metadata = {"source": str(self.file_path)}
        try:
            context = self.extract_content()
            if context:
                yield Document(page_content=context, metadata=metadata)
            else:
                logging.info("No content extracted from PPTX file.")
                yield Document(page_content="", metadata=metadata)
        except Exception as ex:
            logging.exception(f"Error extracting content from file: {self.file_path}.")
            yield Document(page_content="", metadata=metadata)


class ExcelLoader(BaseLoader):
    def __init__(self, file_path: str):
        self.file_path = file_path

    def extract_content(self):
        try:
            if not os.path.exists(self.file_path):
                raise FileNotFoundError(f"Excel file {self.file_path} does not exist.")

            file_ext = os.path.splitext(self.file_path)[-1].lower()

            full_text = ""

            if file_ext == ".xlsx":
                workbook = load_workbook(filename=self.file_path, data_only=True)
                for sheet in workbook.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " ".join(str(cell) for cell in row if cell is not None)
                        if row_text.strip():
                            full_text += row_text.strip() + "\n"
            elif file_ext == ".xls":
                workbook = xlrd.open_workbook(self.file_path)
                for sheet in workbook.sheets():
                    for row_idx in range(sheet.nrows):
                        row_text = " ".join(str(sheet.cell(row_idx, col).value) for col in range(sheet.ncols))
                        if row_text.strip():
                            full_text += row_text.strip() + "\n"
            else:
                raise ValueError(f"Unsupported file format: {file_ext}")

            return full_text.strip()
        except Exception as ex:
            logging.exception("Excel extraction failed.")
            return None

    def lazy_load(self) -> Iterator[Document]:
        metadata = {"source": str(self.file_path)}
        try:
            context = self.extract_content()
            if context:
                yield Document(page_content=context, metadata=metadata)
            else:
                logging.info("No content extracted from Excel file.")
                yield Document(page_content="", metadata=metadata)
        except Exception as ex:
            logging.exception(f"Error extracting content from file: {self.file_path}")
            yield Document(page_content="", metadata=metadata)


def get_loader(filename: str, file_content_type: str, file_path: str):
    file_ext = filename.split(".")[-1].lower()
    known_type = True

    known_source_ext = [
        "go",
        "py",
        "java",
        "sh",
        "bat",
        "ps1",
        "cmd",
        "js",
        "ts",
        "css",
        "cpp",
        "hpp",
        "h",
        "c",
        "cs",
        "sql",
        "log",
        "ini",
        "pl",
        "pm",
        "r",
        "dart",
        "dockerfile",
        "env",
        "php",
        "hs",
        "hsc",
        "lua",
        "nginxconf",
        "conf",
        "m",
        "mm",
        "plsql",
        "perl",
        "rb",
        "rs",
        "db2",
        "scala",
        "bash",
        "swift",
        "vue",
        "svelte",
        "msg",
    ]

    loader: BaseLoader
    if file_ext == "pdf":
        loader = PdfLoader(file_path)
    elif file_ext == "csv":
        loader = CSVLoader(file_path)
    elif file_ext == "rst":
        loader = UnstructuredRSTLoader(file_path, mode="elements")
    elif file_ext == "xml":
        loader = UnstructuredXMLLoader(file_path)
    elif file_ext in ["htm", "html"]:
        loader = BSHTMLLoader(file_path, open_encoding="unicode_escape")
    elif file_ext == "md":
        loader = MarkdownLoader(file_path)
    elif file_content_type == "application/epub+zip":
        loader = UnstructuredEPubLoader(file_path)
    elif (
        file_content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        and file_ext == "docx"
    ):
        loader = Docx2txtLoader(file_path)
    elif file_content_type == "application/msword" and file_ext == "doc":
        loader = UnstructuredWordDocumentLoader(file_path)
    elif file_content_type == "docx/html":
        loader = DocxHtmlLoader(file_path)
    elif file_content_type in [
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ] or file_ext in ["xls", "xlsx"]:
        # loader = UnstructuredExcelLoader(file_path)
        loader = ExcelLoader(file_path)
    elif (
        file_content_type
        in [
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ]
        or file_ext == "pptx"
    ):
        loader = PPTLoader(file_path)
    elif file_ext == "msg":
        loader = OutlookMessageLoader(file_path)
    elif file_ext in known_source_ext or (file_content_type and file_content_type.find("text/") >= 0):
        encoding = fast_detect(file_path)
        if encoding:
            loader = TextLoader(file_path, encoding=encoding)
        else:
            loader = TextLoader(file_path, autodetect_encoding=True)
    else:
        encoding = fast_detect(file_path)
        if encoding:
            loader = TextLoader(file_path, encoding=encoding)
        else:
            loader = TextLoader(file_path, autodetect_encoding=True)
        known_type = False

    return loader, known_type


def fast_detect(file_path: str):
    try:
        rawdata = Path(file_path).read_bytes()
        res = cchardet.detect(rawdata)
        encoding = res.get("encoding", None)
        return encoding
    except Exception as e:
        logging.exception(f"Failed to detect encoding for '{file_path}'.")
        return None


def get_docs(
    file_path: str, file_type: str, chunk_size: int, chunk_overlap: int
) -> tuple[list[Document], list[Document], str]:
    file_name = os.path.basename(file_path)
    loader, known_type = get_loader(file_name, file_type, file_path)
    data = loader.load()
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        add_start_index=True,
    )
    docs = text_splitter.split_documents(data)
    return data, docs, known_type


def random_ua() -> str:
    agent_list = [
        "Mozilla/5.0 (Windows NT 6.1; WOW64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/39.0.2171.95 Safari/537.36 \
        OPR/26.0.1656.60",
        "Opera/8.0 (Windows NT 5.1; U; en)",
        "Mozilla/5.0 (Windows NT 5.1; U; en; rv:1.8.1) Gecko/20061208 Firefox/2.0.0 Opera 9.50",
        "Mozilla/4.0 (compatible; MSIE 6.0; Windows NT 5.1; en) Opera 9.50",
        "Opera/9.80 (Macintosh; Intel Mac OS X 10.6.8; U; en) Presto/2.8.131 Version/11.11",
        "Opera/9.80 (Windows NT 6.1; U; en) Presto/2.8.131 Version/11.11",
        "Opera/9.80 (Android 2.3.4; Linux; Opera Mobi/build-1107180945; U; en-GB) Presto/2.8.149 Version/11.10",
        "Mozilla/5.0 (Windows NT 6.1; WOW64; rv:34.0) Gecko/20100101 Firefox/34.0",
        "Mozilla/5.0 (X11; U; Linux x86_64; zh-CN; rv:1.9.2.10) Gecko/20100922 Ubuntu/10.10 (maverick) Firefox/3.6.10",
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10.6; rv,2.0.1) Gecko/20100101 Firefox/4.0.1",
        "Mozilla/5.0 (Windows NT 6.1; rv,2.0.1) Gecko/20100101 Firefox/4.0.1",
        "Mozilla/5.0 (Windows NT 6.1; WOW64) AppleWebKit/534.57.2 (KHTML, like Gecko) Version/5.1.7 Safari/534.57.2",
        "MAC：Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/81.0.4044.122 \
        Safari/537.36",
        "Windows：Mozilla/5.0 (Windows; U; Windows NT 6.1; en-us) AppleWebKit/534.50 (KHTML, like Gecko) Version/5.1 \
        Safari/534.50",
        "Mozilla/5.0 (iPhone; U; CPU iPhone OS 4_3_3 like Mac OS X; en-us) AppleWebKit/533.17.9 (KHTML, like Gecko) \
        Version/5.0.2 Mobile/8J2 Safari/6533.18.5",
        "Mozilla/5.0 (iPhone; U; CPU iPhone OS 4_3_3 like Mac OS X; en-us) AppleWebKit/533.17.9 (KHTML, like Gecko) \
        Version/5.0.2 Mobile/8J2 Safari/6533.18.5",
        "Mozilla/5.0 (iPad; U; CPU OS 4_3_3 like Mac OS X; en-us) AppleWebKit/533.17.9 (KHTML, like Gecko) \
        Version/5.0.2 Mobile/8J2 Safari/6533.18.5",
        "Mozilla/5.0 (Windows NT 6.1; WOW64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/39.0.2171.71 Safari/537.36",
        "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.11 (KHTML, like Gecko) Chrome/23.0.1271.64 Safari/537.11",
        "Mozilla/5.0 (Windows; U; Windows NT 6.1; en-US) AppleWebKit/534.16 (KHTML, like Gecko) Chrome/10.0.648.133 \
        Safari/534.16",
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_7_0) AppleWebKit/535.11 (KHTML, like Gecko) Chrome/17.0.963.56 \
        Safari/535.11",
        "Mozilla/5.0 (Windows NT 6.1; WOW64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/30.0.1599.101 Safari/537.36",
        "Mozilla/5.0 (Windows NT 6.1; WOW64; Trident/7.0; rv:11.0) like Gecko",
        "Mozilla/4.0 (compatible; MSIE 7.0; Windows NT 5.1; 360SE)",
        "Mozilla/5.0 (Windows NT 6.1; WOW64) AppleWebKit/536.11 (KHTML, like Gecko) Chrome/20.0.1132.11 \
        TaoBrowser/2.0 Safari/536.11",
        "Mozilla/5.0 (Windows NT 6.1; WOW64) AppleWebKit/537.1 (KHTML, like Gecko) Chrome/21.0.1180.71 Safari/537.1 \
        LBBROWSER",
        "Mozilla/5.0 (compatible; MSIE 9.0; Windows NT 6.1; WOW64; Trident/5.0; SLCC2; .NET CLR 2.0.50727; .NET CLR \
        3.5.30729; .NET CLR 3.0.30729; Media Center PC 6.0; .NET4.0C; .NET4.0E; LBBROWSER)",
        "Mozilla/4.0 (compatible; MSIE 6.0; Windows NT 5.1; SV1; QQDownload 732; .NET4.0C; .NET4.0E; LBBROWSER)"
        "Mozilla/5.0 (compatible; MSIE 9.0; Windows NT 6.1; WOW64; Trident/5.0; SLCC2; .NET CLR 2.0.50727; .NET CLR \
        3.5.30729; .NET CLR 3.0.30729; Media Center PC 6.0; .NET4.0C; .NET4.0E; QQBrowser/7.0.3698.400)",
        "Mozilla/4.0 (compatible; MSIE 6.0; Windows NT 5.1; SV1; QQDownload 732; .NET4.0C; .NET4.0E)",
        "Mozilla/5.0 (Windows NT 5.1) AppleWebKit/535.11 (KHTML, like Gecko) Chrome/17.0.963.84 Safari/535.11 SE 2.X \
        MetaSr 1.0",
        "Mozilla/4.0 (compatible; MSIE 7.0; Windows NT 5.1; Trident/4.0; SV1; QQDownload 732; .NET4.0C; .NET4.0E; SE \
        2.X MetaSr 1.0)",
        "Mozilla/4.0 (compatible; MSIE 7.0; Windows NT 5.1; Trident/4.0; SE 2.X MetaSr 1.0; SE 2.X MetaSr 1.0; .NET \
        CLR 2.0.50727; SE 2.X MetaSr 1.0)",
        "Mozilla/5.0 (Windows NT 6.1; WOW64) AppleWebKit/537.36 (KHTML, like Gecko) Maxthon/4.4.3.4000 \
        Chrome/30.0.1599.101 Safari/537.36",
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_7_0) AppleWebKit/535.11 (KHTML, like Gecko) Chrome/17.0.963.56 \
        Safari/535.11",
        "Mozilla/5.0 (Windows NT 6.1; WOW64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/38.0.2125.122 \
        UBrowser/4.0.3214.0 Safari/537.36",
        "Mozilla/5.0 (Windows NT 6.1; WOW64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/55.0.2883.87 \
        UBrowser/6.2.4094.1 Safari/537.36",
        "Mozilla/5.0 (iPhone; U; CPU iPhone OS 4_3_3 like Mac OS X; en-us) AppleWebKit/533.17.9 (KHTML, like Gecko) \
        Version/5.0.2 Mobile/8J2 Safari/6533.18.5",
        "Mozilla/5.0 (iPod; U; CPU iPhone OS 4_3_3 like Mac OS X; en-us) AppleWebKit/533.17.9 (KHTML, like Gecko) \
        Version/5.0.2 Mobile/8J2 Safari/6533.18.5",
        "Mozilla/5.0 (iPad; U; CPU OS 4_2_1 like Mac OS X; zh-cn) AppleWebKit/533.17.9 (KHTML, like Gecko) \
        Version/5.0.2 Mobile/8C148 Safari/6533.18.5",
        "Mozilla/5.0 (iPad; U; CPU OS 4_3_3 like Mac OS X; en-us) AppleWebKit/533.17.9 (KHTML, like Gecko) \
        Version/5.0.2 Mobile/8J2 Safari/6533.18.5",
        "Mozilla/5.0 (Linux; U; Android 2.2.1; zh-cn; HTC_Wildfire_A3333 Build/FRG83D) AppleWebKit/533.1 (KHTML, \
        like Gecko) Version/4.0 Mobile Safari/533.1",
        "Mozilla/5.0 (Linux; U; Android 2.3.7; en-us; Nexus One Build/FRF91) AppleWebKit/533.1 (KHTML, like Gecko) \
        Version/4.0 Mobile Safari/533.1",
        "MQQBrowser/26 Mozilla/5.0 (Linux; U; Android 2.3.7; zh-cn; MB200 Build/GRJ22; CyanogenMod-7) \
        AppleWebKit/533.1 (KHTML, like Gecko) Version/4.0 Mobile Safari/533.1",
        "Opera/9.80 (Android 2.3.4; Linux; Opera Mobi/build-1107180945; U; en-GB) Presto/2.8.149 Version/11.10",
        "Mozilla/5.0 (Linux; U; Android 3.0; en-us; Xoom Build/HRI39) AppleWebKit/534.13 (KHTML, like Gecko) \
        Version/4.0 Safari/534.13",
        "Mozilla/5.0 (BlackBerry; U; BlackBerry 9800; en) AppleWebKit/534.1+ (KHTML, like Gecko) Version/6.0.0.337 \
        Mobile Safari/534.1+",
        "Mozilla/5.0 (hp-tablet; Linux; hpwOS/3.0.0; U; en-US) AppleWebKit/534.6 (KHTML, like Gecko) \
        wOSBrowser/233.70 Safari/534.6 TouchPad/1.0",
        "Mozilla/5.0 (compatible; MSIE 9.0; Windows NT 6.1; Trident/5.0;",
        "Mozilla/4.0 (compatible; MSIE 7.0; Windows NT 6.0)",
        "Mozilla/4.0 (compatible; MSIE 8.0; Windows NT 6.0; Trident/4.0)",
        "Mozilla/4.0 (compatible; MSIE 6.0; Windows NT 5.1)",
        "Mozilla/4.0 (compatible; MSIE 7.0; Windows NT 5.1)",
        "Mozilla/4.0 (compatible; MSIE 7.0; Windows NT 5.1; The World)",
        "Mozilla/4.0 (compatible; MSIE 7.0; Windows NT 5.1; TencentTraveler 4.0)",
        "Mozilla/4.0 (compatible; MSIE 7.0; Windows NT 5.1; Avant Browser)",
        "Mozilla/5.0 (Linux; U; Android 2.3.7; en-us; Nexus One Build/FRF91) AppleWebKit/533.1 (KHTML, like Gecko) \
        Version/4.0 Mobile Safari/533.1",
        "Mozilla/5.0 (SymbianOS/9.4; Series60/5.0 NokiaN97-1/20.0.019; Profile/MIDP-2.1 Configuration/CLDC-1.1) \
        AppleWebKit/525 (KHTML, like Gecko) BrowserNG/7.1.18124",
        "Mozilla/5.0 (compatible; MSIE 9.0; Windows Phone OS 7.5; Trident/5.0; IEMobile/9.0; HTC; Titan)",
    ]
    return agent_list[random.randint(0, len(agent_list) - 1)]


def extract_urls_and_text(text):
    url_pattern = re.compile(r"(https?://(?:www\.)?[-\w]+(?:\.\w[-\w]*)+" r"(?:/[-\w@:%_\+.~#?&//=]*)?)")
    urls = re.findall(url_pattern, text)
    cleaned_text = re.sub(url_pattern, "", text).strip()
    return urls, cleaned_text


def generate_file_abstract(bot_model_config: BotModelConfig, content: str):
    model_dict = bot_model_config.model_dict
    if model_dict is None:
        return ""

    model_name = model_dict.get("name")
    if not model_name or not ollama_model_is_generation(model_name):
        return ""

    provider = model_dict.get("provider", "")

    params = {
        "seed": 101,
        "temperature": 0,
        "top_p": 0.9,
        "top_k": 40,
        "repeat_penalty": 1.2,
        "num_predict": -1,
        "mirostat": 0,
        "mirostat_tau": 5,
        "mirostat_eta": 0.1,
        "repeat_last_n": 64,
        "tfs_z": 1,
        "num_ctx": 32000,
    }
    llm_instance = model_provider_manager.get_model_instance(provider, model_name, params)
    llm_instance.callbacks = [LoggingOutCallbackHandler()]

    prompt = "您是一个有用的助手，负责生成下面提供的数据的全面摘要，确保使用第三人称撰写，摘要尽量简练，\
    最好控制在20字以内"
    prompt_messages = [SystemMessage(content=prompt), HumanMessage(content=content)]

    text = ""
    total_chars = 0
    for chunk in llm_instance.stream(prompt_messages):
        text += chunk.content
        total_chars += len(chunk.content)
        if total_chars >= 200:
            break

    return text.strip()


def get_file_list(dir_name: str) -> dict[str, str]:
    file_dict = {}
    for dir_path, dir_names, file_names in os.walk(dir_name):
        for file_name in file_names:
            if "." not in file_name:
                continue
            if file_name.startswith("."):
                continue
            ext = os.path.splitext(file_name)[1].lower()
            if ext not in [
                ".txt",
                ".docx",
                ".xlsx",
                ".xls",
                ".csv",
                ".pptx",
                ".ppt",
                ".pdf",
                ".md",
                ".json",
                ".html",
            ]:
                continue
            cur_file_path = os.sep.join([dir_path, file_name])
            sha256_hash = hashlib.sha256()
            if os.path.isfile(cur_file_path):
                with open(cur_file_path, "rb") as fp:
                    for byte_block in iter(lambda: fp.read(4096), b""):
                        sha256_hash.update(byte_block)
            file_hash = sha256_hash.hexdigest()
            file_dict[file_hash] = cur_file_path
            if len(file_dict) > 1000:
                return file_dict
        for inner_dir_name in dir_names:
            if len(file_dict) > 1000:
                return file_dict
            file_dict.update(get_file_list(os.sep.join([dir_path, inner_dir_name])))
    return file_dict
